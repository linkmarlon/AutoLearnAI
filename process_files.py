import os
import chromadb
from sentence_transformers import SentenceTransformer
import pdfplumber
from docx import Document
from pptx import Presentation
import zipfile
import rarfile
import py7zr
import tarfile
import gzip
import bz2
import requests
from io import BytesIO
import mega
from gdown import download
import hashlib
from security import validate_url, validate_file, encrypt_file, decrypt_file
import logging
from urllib.parse import urlparse
import json

logging.basicConfig(
    filename='data/errors.log',
    level=logging.ERROR,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger('ProcessFiles')

client = chromadb.Client()
collection = None

def get_file_hash(data):
    try:
        if not isinstance(data, bytes):
            data = bytes(data)
        hasher = hashlib.md5()
        hasher.update(data)
        return hasher.hexdigest()
    except Exception as e:
        logger.error(f"Erro ao calcular hash: {str(e)}")
        raise

def detect_source(url):
    try:
        parsed = urlparse(url)
        netloc = parsed.netloc.lower()
        if 'terabox.com' in netloc or '1024terabox.com' in netloc:
            return "Terabox"
        elif 'mega.nz' in netloc:
            return "MEGA"
        elif 'drive.google.com' in netloc:
            return "Google Drive"
        elif 'docs.google.com' in netloc:
            return "Google Docs"
        logger.warning(f"Fonte desconhecida: {url}")
        return None
    except Exception as e:
        logger.error(f"Erro ao detectar fonte de {url}: {str(e)}")
        return None

def download_from_terabox(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.content
        if not isinstance(content, bytes):
            content = bytes(content)
        return [content]
    except Exception as e:
        logger.error(f"Erro ao baixar do Terabox: {str(e)}")
        raise

def download_from_mega(url):
    try:
        m = mega.Mega().login()
        content = m.download_url(url)
        if not isinstance(content, bytes):
            content = bytes(content)
        return [content]
    except Exception as e:
        logger.error(f"Erro ao baixar do MEGA: {str(e)}")
        raise

def download_from_google_drive(url):
    try:
        file_id = url.split('/d/')[1].split('/')[0]
        response = requests.get(f"https://drive.google.com/uc?export=download&id={file_id}", timeout=10)
        response.raise_for_status()
        content = response.content
        if not isinstance(content, bytes):
            content = bytes(content)
        return [content]
    except Exception as e:
        logger.error(f"Erro ao baixar do Google Drive: {str(e)}")
        raise

def download_from_google_docs(url):
    try:
        file_id = url.split('/d/')[1].split('/')[0]
        response = requests.get(f"https://docs.google.com/document/d/{file_id}/export?format=txt", timeout=10)
        response.raise_for_status()
        content = response.content
        if not isinstance(content, bytes):
            content = bytes(content)
        return [content]
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 401:
            logger.error(f"Google Docs privado, precisa de link público: {url}")
            return []
        raise
    except Exception as e:
        logger.error(f"Erro ao baixar do Google Docs: {str(e)}")
        raise

def download_file(source, url, encryption_key):
    try:
        temp_dir = 'data/'
        os.makedirs(temp_dir, exist_ok=True)
        files_data = []
        
        detected_source = detect_source(url) or source
        if detected_source == "Terabox":
            files_data = download_from_terabox(url)
        elif detected_source == "MEGA":
            files_data = download_from_mega(url)
        elif detected_source == "Google Drive":
            files_data = download_from_google_drive(url)
        elif detected_source == "Google Docs":
            files_data = download_from_google_docs(url)
        else:
            raise ValueError(f"Fonte não suportada para {url}")

        file_paths = []
        for i, data in enumerate(files_data):
            if not isinstance(data, bytes):
                data = bytes(data)
            if validate_file(data):
                key = base64.b64decode(encryption_key)
                if not isinstance(key, bytes):
                    key = bytes(key)
                encrypted_data = encrypt_file(data, key)
                file_path = os.path.join(temp_dir, f'temp_file_{i}.enc')
                with open(file_path, 'w') as f:
                    json.dump(encrypted_data, f)
                file_paths.append(file_path)
            else:
                logger.error(f"Arquivo inválido bloqueado: {url}")
        return file_paths
    except Exception as e:
        logger.error(f"Erro ao baixar arquivo de {detected_source or source}: {str(e)}")
        raise

def extract_text(file_path, encryption_key):
    try:
        with open(file_path, 'r') as f:
            encrypted_data = json.load(f)
        key = base64.b64decode(encryption_key)
        if not isinstance(key, bytes):
            key = bytes(key)
        data = decrypt_file(encrypted_data, key)
        if not data:
            return "[Erro ao descriptografar]"
        temp_file = 'data/temp_decrypted'
        with open(temp_file, 'wb') as f:
            f.write(data)
    except Exception as e:
        logger.error(f"Erro ao preparar arquivo: {str(e)}")
        return "[Erro ao processar arquivo]"

    ext = os.path.splitext(temp_file)[1].lower()
    try:
        if ext == '.pdf':
            with pdfplumber.open(temp_file) as pdf:
                return ''.join(page.extract_text() or '' for page in pdf.pages)
        elif ext == '.docx':
            doc = Document(temp_file)
            return '\n'.join(p.text for p in doc.paragraphs)
        elif ext == '.pptx':
            prs = Presentation(temp_file)
            return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']:
            temp_extract = 'data/extracted/'
            os.makedirs(temp_extract, exist_ok=True)
            if ext == '.zip':
                with zipfile.ZipFile(temp_file, 'r') as z:
                    z.extractall(temp_extract)
            elif ext == '.rar' and rarfile:
                with rarfile.RarFile(temp_file, 'r') as r:
                    r.extractall(temp_extract)
            elif ext == '.7z':
                with py7zr.SevenZipFile(temp_file, 'r') as z:
                    z.extractall(temp_extract)
            elif ext == '.tar':
                with tarfile.open(temp_file, 'r') as t:
                    t.extractall(temp_extract)
            elif ext == '.gz':
                with gzip.open(temp_file, 'rb') as g, open(os.path.join(temp_extract, 'file'), 'wb') as f:
                    f.write(g.read())
            elif ext == '.bz2':
                with bz2.open(temp_file, 'rb') as b, open(os.path.join(temp_extract, 'file'), 'wb') as f:
                    f.write(b.read())
            texts = []
            for root, _, files in os.walk(temp_extract):
                for f in files:
                    file_path = os.path.join(root, f)
                    with open(file_path, 'rb') as ff:
                        temp_data = ff.read()
                    if validate_file(temp_data):
                        key = base64.b64decode(encryption_key)
                        if not isinstance(key, bytes):
                            key = bytes(key)
                        temp_encrypted = encrypt_file(temp_data, key)
                        temp_enc_path = os.path.join(temp_extract, f'enc_{f}.enc')
                        with open(temp_enc_path, 'w') as ff:
                            json.dump(temp_encrypted, ff)
                        texts.append(extract_text(temp_enc_path, encryption_key))
                    else:
                        logger.error(f"Arquivo extraído bloqueado: {f}")
            return '\n'.join(texts)
        else:
            try:
                with open(temp_file, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except:
                return "[Arquivo binário, processado como dados brutos]"
    except Exception as e:
        logger.error(f"Erro ao extrair texto de {file_path}: {str(e)}")
        return f"[Erro ao ler arquivo: {e}]"
    finally:
        if os.path.exists(temp_file):
            os.remove(temp_file)

def process_files(files, source='Computador', shared=False, encryption_key=None):
    global collection
    texts = []
    file_hashes = []
    try:
        if shared and collection:
            existing_hashes = collection.get()['metadatas'] or []
            existing_hashes = [m.get('hash', '') for m in existing_hashes if m]
        else:
            existing_hashes = []

        if source == 'Computador':
            os.makedirs('data/', exist_ok=True)
            for file in files:
                try:
                    file_data = file.read()
                    if not isinstance(file_data, bytes):
                        file_data = bytes(file_data)
                    if validate_file(file_data):
                        file_hash = get_file_hash(file_data)
                        if file_hash not in existing_hashes:
                            key = base64.b64decode(encryption_key)
                            if not isinstance(key, bytes):
                                key = bytes(key)
                            encrypted_data = encrypt_file(file_data, key)
                            file_path = os.path.join('data/', file.name + '.enc')
                            with open(file_path, 'w') as f:
                                json.dump(encrypted_data, f)
                            texts.append(extract_text(file_path, encryption_key))
                            file_hashes.append(file_hash)
                        else:
                            logger.info(f"Arquivo duplicado: {file.name}")
                    else:
                        logger.error(f"Arquivo inválido: {file.name}")
                except Exception as e:
                    logger.error(f"Erro ao processar arquivo local {file.name}: {str(e)}")
        else:
            for url in files:
                if url.strip():
                    try:
                        file_paths = download_file(source, url, encryption_key)
                        for file_path in file_paths:
                            with open(file_path, 'r') as f:
                                encrypted_data = json.load(f)
                            key = base64.b64decode(encryption_key)
                            if not isinstance(key, bytes):
                                key = bytes(key)
                            decrypted_data = decrypt_file(encrypted_data, key)
                            if decrypted_data:
                                if not isinstance(decrypted_data, bytes):
                                    decrypted_data = bytes(decrypted_data)
                                file_hash = get_file_hash(decrypted_data)
                                if file_hash not in existing_hashes:
                                    texts.append(extract_text(file_path, encryption_key))
                                    file_hashes.append(file_hash)
                            else:
                                logger.error(f"Falha ao descriptografar: {url}")
                    except Exception as e:
                        logger.error(f"Erro ao processar URL {url}: {str(e)}")

        if texts:
            model = SentenceTransformer('all-MiniLM-L6-v2')
            embeddings = model.encode(texts)
            if collection is None:
                collection = client.create_collection('autolearnai')
            collection.add(
                documents=texts,
                embeddings=embeddings,
                ids=[f'doc_{i}' for i in range(len(texts))],
                metadatas=[{'hash': h} for h in file_hashes]
            )
        return collection
    except Exception as e:
        logger.critical(f"Erro fatal em process_files: {str(e)}")
        raise

def get_collection():
    global collection
    try:
        if collection is None:
            try:
                collection = client.get_collection('autolearnai')
            except:
                collection = client.create_collection('autolearnai')
        return collection
    except Exception as e:
        logger.error(f"Erro ao acessar coleção: {str(e)}")
        raise